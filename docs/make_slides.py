# -*- coding: utf-8 -*-
"""توليد عرض المناقشة النهائية لمشروع عطريقك — هوية بصرية موحدة RTL."""
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─── هوية المشروع (من lib/core/them/my_colors.dart) ───
NAVY = RGBColor(0x25, 0x3C, 0x5B)      # primary
DARK_NAVY = RGBColor(0x0D, 0x1B, 0x2A)  # navy
ORANGE = RGBColor(0xED, 0x8B, 0x10)     # accent
BLUE = RGBColor(0x35, 0x68, 0x99)
CYAN = RGBColor(0x5F, 0xBB, 0xC8)
BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x26, 0x32, 0x38)
TEXT2 = RGBColor(0x50, 0x5A, 0x63)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

FONT = "Segoe UI"
LOGO = r"d:\my_projects\flutter\alatarekak\assets\images\app_logo.png"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


_LATIN = re.compile(r"([A-Za-z][A-Za-z0-9./+%_-]*(?: [A-Za-z][A-Za-z0-9./+%_-]*)*)")


def bidi(text):
    """عزل المقاطع اللاتينية (مع أقواسها) كوحدات LTR داخل النص العربي.

    الأقواس ذات المحتوى العربي تُترك كما هي — خوارزمية bidi تعكسها صحيحاً.
    """
    holes = []

    def stash(m):
        holes.append(m.group(0))
        return "\x00{}\x01".format(len(holes) - 1)

    # أقواس تحيط بمحتوى لاتيني بالكامل: تُحفظ ثم تُعزل كوحدة واحدة
    text = re.sub(r"\([A-Za-z0-9 ./+%_-]+\)", stash, text)
    text = _LATIN.sub("‎\\1‏", text)
    text = text.replace("—", "‏—‏")
    for i, h in enumerate(holes):
        text = text.replace("\x00{}\x01".format(i), "‎" + h + "‏")
    return text


def set_text(run, txt):
    run.text = bidi(txt)


def style_run(run, size, color, bold=False, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    cs = rPr.find(qn("a:cs"))
    if cs is None:
        cs = rPr.makeelement(qn("a:cs"), {})
        rPr.append(cs)
    cs.set("typeface", name)


def rtl(paragraph, align=PP_ALIGN.RIGHT):
    paragraph.alignment = align
    paragraph._p.get_or_add_pPr().set("rtl", "1")


def add_rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, bold) أو list of runs-lists."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        runs = line if isinstance(line, list) else [line]
        for (txt, size, color, bold) in runs:
            r = p.add_run()
            set_text(r, txt)
            style_run(r, size, color, bold)
        rtl(p, align)
        p.space_after = Pt(6)
    return tb


def add_bullets(slide, x, y, w, h, items, size=17, gap=10):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        head, _, rest = item.partition("|")
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = "●  "
        style_run(r, size - 4, ORANGE, True)
        r = p.add_run()
        set_text(r, head.strip())
        style_run(r, size, TEXT, bool(rest))
        if rest:
            r2 = p.add_run()
            set_text(r2, "  —  " + rest.strip())
            style_run(r2, size - 1.5, TEXT2, False)
        rtl(p)
        p.space_after = Pt(gap)
    return tb


def logo_badge(slide, cx, cy, d):
    """لوغو داخل دائرة بيضاء (السهام البيضاء داخل اللوغو تبقى ظاهرة)."""
    add_rect(slide, cx, cy, d, d, WHITE, MSO_SHAPE.OVAL)
    pad = Emu(int(d * 0.12))
    slide.shapes.add_picture(LOGO, cx + pad, cy + pad,
                             width=d - 2 * pad, height=d - 2 * pad)


def footer(slide, n):
    add_rect(slide, 0, SH - Pt(4), SW, Pt(4), ORANGE)
    tb = add_text(slide, Inches(0.25), SH - Inches(0.42), Inches(1.2),
                  Inches(0.3), [(str(n), 11, TEXT2, False)],
                  align=PP_ALIGN.LEFT)
    slide.shapes.add_picture(LOGO, SW - Inches(0.55), SH - Inches(0.52),
                             height=Inches(0.38))


def content_slide(title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, WHITE)
    # شريط جانبي أيمن رفيع بلون الهوية
    add_rect(s, SW - Inches(0.12), 0, Inches(0.12), SH, NAVY)
    # العنوان
    add_text(s, Inches(0.7), Inches(0.35), Inches(11.9), Inches(0.9),
             [(title, 30, NAVY, True)])
    add_rect(s, SW - Inches(3.4), Inches(1.12), Inches(2.7), Pt(3.5), ORANGE)
    if subtitle:
        add_text(s, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5),
                 [(subtitle, 15, TEXT2, False)])
    return s


def card(slide, x, y, w, h, title, body, accent=ORANGE, tsize=16, bsize=12.5):
    add_rect(slide, x, y, w, h, BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(slide, x + w - Inches(0.09), y + Inches(0.15),
             Inches(0.09), h - Inches(0.3), accent)
    add_text(slide, x + Inches(0.15), y + Inches(0.12), w - Inches(0.45),
             Inches(0.5), [(title, tsize, NAVY, True)])
    add_text(slide, x + Inches(0.15), y + Inches(0.62), w - Inches(0.45),
             h - Inches(0.75), [(body, bsize, TEXT2, False)])


N = 0


def num():
    global N
    N += 1
    return N


# ════════════════ 1. الغلاف ════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, DARK_NAVY)
add_rect(s, 0, 0, SW, Pt(6), ORANGE)
add_rect(s, 0, SH - Pt(6), SW, Pt(6), ORANGE)
logo_badge(s, SW / 2 - Inches(1.1), Inches(0.8), Inches(2.2))
add_text(s, Inches(1.5), Inches(3.15), Inches(10.33), Inches(1.2),
         [("عطريقك", 60, WHITE, True)], align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.35), Inches(10.33), Inches(0.6),
         [("منصة مشاركة الرحلات بين المدن — مقعد فارغ لواحد، وصول أوفر للجميع",
           20, ORANGE, False)], align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(5.5), Inches(10.33), Inches(1.2),
         [("مناقشة مشروع التخرج — العرض النهائي", 16, WHITE, False),
          ("تطبيق Flutter  •  واجهة Laravel API  •  لوحة إدارة ويب",
           13, CYAN, False),
          ("2026", 13, WHITE, False)], align=PP_ALIGN.CENTER)

# ════════════════ 2. المشكلة ════════════════
s = content_slide("المشكلة — التنقل بين المدن اليوم", "لماذا يحتاج السوق حلاً؟")
add_bullets(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5), [
    "تكلفة مرتفعة | التاكسي الخاص بين المدن مكلف جداً لفرد واحد، والنقل الجماعي غير مرن في المواعيد والمسارات",
    "مقاعد فارغة مهدورة | آلاف السيارات تتحرك يومياً بين المدن وفيها 3 مقاعد فارغة — طاقة نقل ضائعة",
    "انعدام الثقة | التشارك الحالي يتم عبر مجموعات فيسبوك/واتساب: لا توثيق للهوية، لا تقييمات، لا ضمانات",
    "لا ضمانات مالية | الدفع نقدي بلا إثبات، الإلغاءات المتأخرة بلا تعويض، والخلافات بلا جهة تحكيم",
    "غياب السلامة | لا سجل للسائق أو مركبته، ولا آلية للإبلاغ أو تتبع الرحلة أو مشاركة تفاصيلها",
], size=17, gap=14)

# ════════════════ 3. الحل ════════════════
s = content_slide("الحل — عطريقك", "منصة موثوقة تحول المقاعد الفارغة إلى شبكة نقل")
add_bullets(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(3.2), [
    "سائق موثّق ينشر رحلته | هوية + رخصة + وثيقة مركبة تراجعها الإدارة قبل أول رحلة (BR-01)",
    "راكب يبحث ويحجز بثقة | بحث جغرافي بالإحداثيات والتاريخ وعدد المقاعد، مع ملف كامل للسائق وتقييماته",
    "دفع إلكتروني عبر محفظة داخلية | خصم آمن عند الحجز، واسترداد عادل متدرج عند الإلغاء",
    "منظومة ثقة متكاملة | نقاط ثقة سلوكية، تقييمات وتعليقات، بلاغات عدم حضور، وحظر للمخالفين",
])
row_y = Inches(5.2)
for i, (t, b) in enumerate([
        ("للراكب", "وصول أرخص وأسرع بمرونة مواعيد حقيقية"),
        ("للسائق", "تغطية تكاليف الرحلة من مقاعد كانت ستبقى فارغة"),
        ("للمجتمع", "ازدحام أقل، استهلاك وقود أقل، شبكة معارف أوسع")]):
    card(s, Inches(0.7) + i * Inches(4.1), row_y, Inches(3.8), Inches(1.5), t, b)
footer(s, num() + 1)

# ════════════════ 4. لماذا سيختارنا المستخدم ════════════════
s = content_slide("لماذا يختار المستخدم عطريقك؟",
                  "مقارنة مباشرة مع البدائل المتاحة فعلياً")
rows = [
    ("المعيار", "مجموعات فيسبوك/واتساب", "تاكسي خاص", "عطريقك"),
    ("توثيق هوية السائق", "لا يوجد", "غير معياري", "إلزامي قبل أول رحلة"),
    ("السعر", "متفاوت وغير ملزم", "مرتفع جداً", "سعر معلن ثابت للمقعد"),
    ("ضمان الحجز", "وعود شفهية", "حسب المكتب", "حجز مؤكد بخصم من المحفظة"),
    ("الإلغاء والاسترداد", "خسارة كاملة", "لا سياسة", "استرداد متدرج حسب التوقيت"),
    ("التقييم والمحاسبة", "لا يوجد", "لا يوجد", "نقاط ثقة + تقييمات + حظر"),
    ("التواصل", "أرقام مكشوفة", "هاتف فقط", "شات داخل التطبيق + إشعارات"),
]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.7), Inches(1.7),
                         Inches(11.9), Inches(4.9)).table
tbl.columns[0].width = Inches(3.1)
tbl.columns[1].width = Inches(3.2)
tbl.columns[2].width = Inches(2.4)
tbl.columns[3].width = Inches(3.2)
for ri, row in enumerate(rows):
    # عكس الأعمدة لاتجاه RTL: المعيار في أقصى اليمين
    for ci, val in enumerate(reversed(row)):
        cell = tbl.cell(ri, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        set_text(r, val)
        is_header = ri == 0
        is_us = ci == 0 and ri > 0
        is_crit = ci == 3 and ri > 0
        style_run(r, 13 if not is_header else 14,
                  WHITE if is_header else (GREEN if is_us else (NAVY if is_crit else TEXT2)),
                  is_header or is_us or is_crit)
        rtl(p, PP_ALIGN.CENTER)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if is_header else (WHITE if ri % 2 else BG)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
footer(s, num() + 1)

# ════════════════ 5. نظرة عامة على النظام ════════════════
s = content_slide("بنية النظام — ثلاثة مكونات متكاملة")
comps = [
    ("تطبيق الجوال — Flutter", "واجهة الراكب والسائق: بحث، حجز، إنشاء رحلات،"
     " محفظة، شات، إشعارات — عربي RTL مع وضع داكن", ORANGE),
    ("الخادم — Laravel REST API", "منطق الأعمال وقواعده، JWT + Refresh،"
     " MySQL، بث فوري عبر WebSockets، معالجة المدفوعات والاسترداد", NAVY),
    ("لوحة الإدارة — ويب", "مراجعة طلبات التوثيق، إدارة طلبات شحن/سحب المحفظة،"
     " الشكاوى والدعم، مراقبة المستخدمين والحظر", CYAN),
]
for i, (t, b, c) in enumerate(comps):
    card(s, Inches(0.7) + i * Inches(4.1), Inches(1.9), Inches(3.8),
         Inches(2.5), t, b, accent=c, tsize=16, bsize=13)
add_rect(s, Inches(2.2), Inches(4.75), Inches(8.9), Pt(2.5), NAVY)
add_text(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.6), [
    ("تكامل خارجي:", 16, NAVY, True),
    ("خرائط OpenStreetMap + MapTiler  •  مسارات GraphHopper / OpenRouteService"
     "  •  إشعارات Firebase Cloud Messaging  •  بث Pusher Protocol",
     14, TEXT2, False),
])
footer(s, num() + 1)

# ════════════════ 6. التقنيات ════════════════
s = content_slide("لماذا هذه التقنيات؟", "كل اختيار له مبرر هندسي")
add_bullets(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(5), [
    "Flutter | قاعدة كود واحدة لـ Android و iOS بأداء أصلي — فريق صغير يغطي منصتين",
    "Cubit (flutter_bloc) | فصل كامل بين الواجهة ومنطق الأعمال + قابلية اختبار عالية (أثبتناها بـ 105 اختبارات)",
    "Laravel | إطار ناضج بنظام صلاحيات وMiddleware جاهز، يسرّع بناء REST API آمن",
    "MySQL | بيانات علائقية بطبيعتها (مستخدم/رحلة/حجز/محفظة) مع معاملات ACID للعمليات المالية",
    "Hive | تخزين محلي سريع بلا SQL للكاش والجلسة — مع تشفير AES-256 لصندوق الجلسة",
    "Dio | عميل HTTP قابل للتوسعة بـ Interceptors: توكن تلقائي، تجديد جلسة، إعادة محاولة",
    "OpenStreetMap بدل Google Maps | مجاني بلا مفاتيح فوترة، وكافٍ تماماً لاختيار المواقع والمسارات",
])
footer(s, num() + 1)

# ════════════════ 7. المعمارية ════════════════
s = content_slide("معمارية التطبيق — Clean Architecture",
                  "NFR-24: فصل الطبقات وقابلية الصيانة")
layers = [
    ("Presentation", "الشاشات + Cubits — تستهلك حالات وتعرضها، لا تعرف شيئاً عن الشبكة", ORANGE),
    ("Domain", "الكيانات وعقود الـ Repositories — قلب منطق الأعمال المستقل", BLUE),
    ("Data", "Repositories + Data Sources + Models — Dio للشبكة و Hive للتخزين", NAVY),
]
for i, (t, b, c) in enumerate(layers):
    y = Inches(1.8) + i * Inches(1.35)
    add_rect(s, Inches(3.2), y, Inches(6.9), Inches(1.1), BG,
             MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(s, Inches(9.75), y + Inches(0.1), Inches(0.12), Inches(0.9), c)
    add_text(s, Inches(3.4), y + Inches(0.08), Inches(6.3), Inches(0.45),
             [(t, 16, NAVY, True)])
    add_text(s, Inches(3.4), y + Inches(0.52), Inches(6.3), Inches(0.5),
             [(b, 12, TEXT2, False)])
add_text(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1), [
    [("القاعدة الذهبية:  ", 15, NAVY, True),
     ("الاعتماديات تتجه للداخل فقط — تغيير مكتبة الشبكة أو التخزين لا يلمس منطق الأعمال أو الشاشات. "
      "حقن التبعيات عبر GetIt، وإدارة الحالة عبر Cubit مع كل تدفق يمر بحالات Loading / Success / Error.",
      13.5, TEXT2, False)],
])
footer(s, num() + 1)

# ════════════════ 8. رحلة الراكب ════════════════
s = content_slide("رحلة الراكب — من البحث إلى التقييم")
steps = [
    ("1. بحث", "إحداثيات انطلاق ووجهة + تاريخ + مقاعد"),
    ("2. اختيار", "ملف السائق، تقييمه، سعر المقعد، نوع المركبة"),
    ("3. حجز", "فوري أو بطلب موافقة — خصم آمن من المحفظة"),
    ("4. تواصل", "شات فوري مع السائق + إشعارات لحظية"),
    ("5. الرحلة", "تأكيد ثنائي للإنهاء (سائق ثم راكب)"),
    ("6. تقييم", "نجوم + تعليق يبنيان سمعة الطرفين"),
]
for i, (t, b) in enumerate(steps):
    x = Inches(0.7) + (i % 3) * Inches(4.1)
    y = Inches(1.9) + (i // 3) * Inches(2.3)
    card(s, x, y, Inches(3.8), Inches(2.0), t, b, tsize=17, bsize=13)
add_text(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.6),
         [[("حماية الراكب:  ", 14, NAVY, True),
           ("إلغاء باسترداد متدرج، بلاغ «السائق لم يحضر» باسترداد كامل، وزر دعم مباشر",
            13, TEXT2, False)]])
footer(s, num() + 1)

# ════════════════ 9. رحلة السائق ════════════════
s = content_slide("رحلة السائق — إنشاء الرحلة بمعالج ذكي")
steps = [
    ("التوثيق أولاً", "لا نشر رحلات قبل قبول الإدارة للهوية والرخصة ووثيقة المركبة (BR-01)"),
    ("تحديد المواقع", "اختيار الانطلاق والوجهة من الخريطة مباشرة أو بالبحث الجغرافي"),
    ("مسارات بديلة", "حساب عدة مسارات فعلية (مسافة وزمن) واختيار الأنسب"),
    ("التفاصيل والتسعير", "مقاعد (بحد أقصى 4 — BR-02)، سعر المقعد، نوع الحجز، موعد الانطلاق"),
    ("منع التعارض", "النظام يرفض رحلتين متداخلتين زمنياً لنفس السائق (BR-03)"),
    ("إدارة الركاب", "قبول/رفض طلبات الحجز، بلاغ عدم حضور راكب، إنهاء وتأكيد"),
]
for i, (t, b) in enumerate(steps):
    x = Inches(0.7) + (i % 3) * Inches(4.1)
    y = Inches(1.9) + (i // 3) * Inches(2.4)
    card(s, x, y, Inches(3.8), Inches(2.1), t, b, accent=BLUE, tsize=15.5,
         bsize=12.5)
footer(s, num() + 1)

# ════════════════ 10. المصادقة والجلسات ════════════════
s = content_slide("المصادقة وإدارة الجلسات", "NFR-10 + NFR-14")
add_bullets(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(5), [
    "تسجيل بتحقق OTP عبر البريد | الحساب لا يُفعّل قبل إدخال رمز 6 أرقام — مع مؤقت إعادة إرسال وحد محاولات",
    "JWT + Refresh Token | جلسات بلا حالة على الخادم؛ access قصير العمر و refresh لتجديد صامت",
    "تجديد تلقائي شفاف | Interceptor يلتقط 401، يجدد التوكن، ويعيد الطلب الأصلي دون أن يشعر المستخدم",
    "أكواد جلسة مصنفة | TOKEN_INVALID يجدد، TOKEN_INVALIDATED/USER_INACTIVE خروج فوري، USER_BANNED شاشة حظر بالتفاصيل",
    "تخزين مشفر AES-256 | التوكنات في صندوق Hive مشفر بمفتاح داخل Android Keystore / iOS Keychain — ليست في SharedPreferences إطلاقاً",
    "استعادة كلمة المرور | بريد ← OTP ← رمز إعادة تعيين لمرة واحدة بصلاحية محدودة",
])
footer(s, num() + 1)

# ════════════════ 11. المحفظة والدفع ════════════════
s = content_slide("المحفظة والمدفوعات", "عمليات مالية آمنة داخل المنصة")
add_bullets(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(3.6), [
    "تفعيل بمصادقة مزدوجة | إنشاء المحفظة يتطلب كلمة المرور + OTP على رقم الهاتف",
    "شحن وسحب بطلبات مُدارة | الطلب يُراجع من الإدارة عبر لوحة الويب ثم يُقيد في الرصيد",
    "خصم عند الحجز لحظياً | ولا يكتمل أي خصم إلا بعد استجابة الخادم (SAF-04) — مع مفتاح idempotency يمنع التكرار",
    "استرداد متدرج عادل | نسبة الاسترداد تتناقص كلما اقترب موعد الانطلاق — توازن بين حق الراكب وضرر السائق",
    "استرداد كامل تلقائي | عند إلغاء السائق أو ثبوت عدم حضوره",
], size=16, gap=9)
row_y = Inches(5.35)
for i, (t, b) in enumerate([
        ("سجل معاملات كامل", "كل حركة (شحن/خصم/استرداد) مؤرخة وقابلة للمراجعة"),
        ("حدود وحماية", "حد أقصى للخصم اليومي + رصد النشاط غير المعتاد (SAF-05/06)"),
        ("لماذا محفظة داخلية؟", "لا بوابات دفع متاحة محلياً — المحفظة تمنح تجربة إلكترونية كاملة اليوم وقابلة للربط ببوابة لاحقاً")]):
    card(s, Inches(0.7) + i * Inches(4.1), row_y, Inches(3.8), Inches(1.55),
         t, b, bsize=11.5, tsize=14)
footer(s, num() + 1)

# ════════════════ 12. الشات والإشعارات ════════════════
s = content_slide("التواصل الفوري — شات وإشعارات")
card(s, Inches(6.9), Inches(1.8), Inches(5.7), Inches(4.4),
     "الشات الفوري (WebSockets)",
     "بث لحظي عبر بروتوكول Pusher مع قنوات خاصة لكل محادثة\n\n"
     "رسائل نصية وصور، حالة قراءة، وحذف رسائل\n\n"
     "إعادة اتصال تلقائية عند عودة الشبكة ومزامنة ما فات\n\n"
     "محادثة تُفتح تلقائياً بين الراكب والسائق عند تأكيد الحجز",
     accent=ORANGE, bsize=13.5)
card(s, Inches(0.7), Inches(1.8), Inches(5.7), Inches(4.4),
     "الإشعارات (FCM + مركز إشعارات)",
     "إشعارات Push عبر Firebase حتى والتطبيق مغلق\n\n"
     "مركز إشعارات داخلي: تصنيفات، تعليم كمقروء، حذف، عداد Badge\n\n"
     "تحديث تفاؤلي فوري مع مزامنة خلفية وكاش يعمل بلا اتصال\n\n"
     "ربط عميق: الضغط على الإشعار يفتح المحادثة أو الرحلة المعنية",
     accent=CYAN, bsize=13.5)
footer(s, num() + 1)

# ════════════════ 13. نقاط الثقة ════════════════
s = content_slide("منظومة الثقة — Trust Score", "الردع الذاتي قبل العقوبة")
add_bullets(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(3.4), [
    "رصيد سلوكي لكل مستخدم | يرتفع بإكمال الرحلات ويهبط بالإلغاءات وبلاغات عدم الحضور",
    "عتبات ملزمة | إنشاء رحلة يتطلب 50 نقطة، والحجز 40 — تُطبق في الخادم وتُعطّل الأزرار في التطبيق استباقياً",
    "شفافية كاملة | شاشة تعرض النقاط والفئة (برونزي…بلاتيني) وسجل كل تغيير وسببه",
    "التقييمات والتعليقات | متوسط نجوم علني + تعليقات موثقة من ركاب فعليين فقط",
    "التصعيد | بلاغات متكررة ← تعليق تلقائي أو حظر إداري مؤقت/دائم مع إظهار السبب",
], size=16, gap=11)
add_text(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1),
         [[("لماذا نقاط وليس تقييماً فقط؟  ", 15, NAVY, True),
           ("التقييم رأي شخصي قد يُجامل؛ النقاط تُحتسب من سلوك مُقاس (إلغاءات، حضور، إكمال) فلا يمكن تجميلها.",
            13.5, TEXT2, False)]])
footer(s, num() + 1)

# ════════════════ 14. السلامة والدعم ════════════════
s = content_slide("السلامة والدعم", "متطلبات SAF + قناة دعم مباشرة")
card(s, Inches(6.9), Inches(1.8), Inches(5.7), Inches(4.5),
     "حماية المستخدم",
     "توثيق إلزامي للسائقين قبل أول رحلة (SAF-03)\n\n"
     "بلاغ «لم يحضر» للطرفين مع أثر مالي ونقاط ثقة\n\n"
     "أرقام تواصل تظهر فقط بعد تأكيد الحجز\n\n"
     "حظر مؤقت/دائم بواجهة توضح السبب والمدة\n\n"
     "خارطة الطريق: زر طوارئ ومشاركة الرحلة (SAF-01/02)",
     accent=GREEN, bsize=13)
card(s, Inches(0.7), Inches(1.8), Inches(5.7), Inches(4.5),
     "الدعم والشكاوى",
     "نظام شكاوى بأنواع ومرفقات (حتى 3 ملفات/شكوى)\n\n"
     "متابعة حالة الشكوى: مفتوحة ← قيد المعالجة ← محلولة\n\n"
     "شات مباشر مع فريق الدعم من داخل التطبيق\n\n"
     "إشعار فوري عند رد الإدارة على الشكوى\n\n"
     "كاش محلي: قائمة شكاواك متاحة حتى بلا اتصال",
     accent=BLUE, bsize=13)
footer(s, num() + 1)

# ════════════════ 15. الجودة والاختبارات ════════════════
s = content_slide("الجودة — اختبارات آلية شاملة", "NFR-02")
stats = [("105", "اختباراً آلياً — كلها ناجحة"),
         ("%77.7", "تغطية أسطر منطق الأعمال المختبَر"),
         ("19", "ملف اختبار — وحدات وواجهات وتكامل"),
         ("16", "ثانية لتشغيل المجموعة كاملة")]
for i, (v, t) in enumerate(stats):
    x = Inches(0.7) + i * Inches(3.075)
    add_rect(s, x, Inches(1.75), Inches(2.87), Inches(1.5), DARK_NAVY,
             MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x, Inches(1.85), Inches(2.87), Inches(0.7),
             [(v, 30, ORANGE, True)], align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), Inches(2.55), Inches(2.67), Inches(0.7),
             [(t, 11.5, WHITE, False)], align=PP_ALIGN.CENTER)
add_bullets(s, Inches(0.7), Inches(3.6), Inches(11.9), Inches(3.2), [
    "Unit Tests | 12 وحدة حالة (مصادقة، حجز، رحلات، محفظة، نقاط، إشعارات) بمحاكاة كاملة للخادم — حالات نجاح وفشل معرّبة",
    "Widget Tests | مكونات مشتركة وشاشات كاملة بكيوبت وهمي دون أي شبكة",
    "Integration Test | إقلاع التطبيق كاملاً على جهاز حقيقي والتحقق من العربية وRTL",
    "اختبارات أمان فعلية | قراءة بايتات القرص لإثبات أن التوكن لم يعد نصاً مكشوفاً",
    "اختبارات الشبكة | إثبات إعادة محاولة GET وعدم تكرار POST مطلقاً (حماية العمليات المالية)",
], size=14.5, gap=8)
footer(s, num() + 1)

# ════════════════ 16. الأمان ════════════════
s = content_slide("الأمان — دفاع متعدد الطبقات")
rows = [
    ("الطبقة", "الإجراء المطبق"),
    ("النقل", "HTTPS/TLS لكل الاتصالات مع الخادم (NFR-12/15)"),
    ("الجلسات", "JWT قصير العمر + Refresh Token + إبطال فوري عند الحظر (NFR-10)"),
    ("التخزين على الجهاز", "AES-256 لصندوق الجلسة، المفتاح في Keystore/Keychain (NFR-14)"),
    ("الخادم", "تحديد معدل الطلبات 60/دقيقة و OTP بمحاولات محدودة (NFR-11) + CSRF/XSS للوحة الويب (NFR-08/09)"),
    ("العمليات المالية", "idempotency key + تأكيد الخادم قبل أي خصم + سقف يومي (SAF-04/05)"),
    ("الحساب", "رصد النشاط غير المعتاد وتعليق تلقائي + حظر إداري (SAF-06)"),
]
tbl = s.shapes.add_table(len(rows), 2, Inches(0.7), Inches(1.75),
                         Inches(11.9), Inches(4.7)).table
tbl.columns[0].width = Inches(8.9)
tbl.columns[1].width = Inches(3.0)
for ri, (layer, action) in enumerate(rows):
    for ci, val in enumerate([action, layer]):
        cell = tbl.cell(ri, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        set_text(r, val)
        head = ri == 0
        style_run(r, 13.5 if not head else 14.5,
                  WHITE if head else (NAVY if ci == 1 else TEXT2),
                  head or ci == 1)
        rtl(p, PP_ALIGN.CENTER if ci == 1 else PP_ALIGN.RIGHT)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if head else (WHITE if ri % 2 else BG)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
footer(s, num() + 1)

# ════════════════ 17. الموثوقية ════════════════
s = content_slide("الموثوقية — التطبيق يعمل على شبكات سيئة", "NFR-17")
add_bullets(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(5), [
    "إعادة محاولة ذكية | أخطاء الشبكة العابرة تُعاد تلقائياً بتراجع تصاعدي — لطلبات القراءة فقط؛ الحجز والدفع لا يُكرران أبداً",
    "مهلات صريحة | 15 ثانية اتصال / 25 استلام — لا طلب يعلق للأبد، ورسالة عربية واضحة عند التجاوز",
    "شريط انقطاع عام | يظهر فوق كل الشاشات لحظة فقدان الشبكة ويختفي عند عودتها",
    "كاش «اعرض ثم حدّث» | النقاط والإشعارات والشكاوى تُعرض من الكاش فوراً ثم تتحدث من الشبكة بصمت",
    "تجديد جلسة بطابور | QueuedInterceptor يمنع تسابق طلبات متزامنة على تجديد التوكن نفسه",
    "إعادة اتصال الشات | السوكيت يعيد الاتصال تلقائياً ويزامن الرسائل الفائتة",
    "رسائل خطأ معرّبة مركزياً | كل رسائل الخادم تمر بمترجم واحد — لا يرى المستخدم نصاً إنجليزياً أو تقنياً",
])
footer(s, num() + 1)

# ════════════════ 18. تجربة الاستخدام ════════════════
s = content_slide("تجربة الاستخدام", "NFR-21 / 23 / 26")
cards = [
    ("عربي RTL بالكامل", "كل الواجهات والرسائل والتواريخ عربية، بترميز UTF-8 واتجاه صحيح"),
    ("وضع داكن كامل", "لوحتا ألوان بهوية واحدة (كحلي/برتقالي) مع تبديل فوري وحفظ التفضيل"),
    ("حالات واضحة دائماً", "لكل شاشة Loading / Success / Error / Empty — لا شاشات بيضاء صامتة"),
    ("تكيف مع كل المقاسات", "ScreenUtil بمقاس تصميم موحد 375×812 على Android و iOS"),
    ("تحديث تفاؤلي", "الإجراءات تظهر نتيجتها فوراً وتُزامن خلفياً — إحساس بالسرعة"),
    ("سحب للتحديث وترقيم صفحات", "قوائم طويلة بلا انتظار: تحميل تدريجي عند التمرير"),
]
for i, (t, b) in enumerate(cards):
    x = Inches(0.7) + (i % 3) * Inches(4.1)
    y = Inches(1.85) + (i // 3) * Inches(2.35)
    card(s, x, y, Inches(3.8), Inches(2.05), t, b, tsize=15, bsize=12.5)
footer(s, num() + 1)

# ════════════════ 19. تحديات وحلول ════════════════
s = content_slide("تحديات واجهتنا — وكيف عالجناها")
add_bullets(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(5), [
    "أشكال أخطاء متعددة من الخادم | بنينا كائن فشل موحداً يفهم 3 مغلفات مختلفة + مترجماً مركزياً للعربية — مع اختبارات لكل شكل",
    "تعارض تجديد التوكن المتزامن | QueuedInterceptor يصفّف الطلبات أثناء التجديد فلا يُطرد المستخدم بسبب سباق طلبات",
    "تشفير التخزين دون فقدان الجلسات | ترحيل تلقائي لمرة واحدة من الصندوق المكشوف إلى المشفر — صفر خروج إجباري",
    "منع ازدواج العمليات المالية | مفتاح idempotency لكل حجز + سياسة «لا إعادة تلقائية لأي طلب كتابة»",
    "الشات الفوري على شبكة متقلبة | إعادة اتصال ذاتية + مزامنة الرسائل الفائتة + إرسال عبر REST والاستقبال عبر السوكيت",
    "عدالة الإلغاء | سياسة استرداد متدرجة زمنياً بدل قرار ثنائي — درسناها كقاعدة عمل قبل برمجتها",
])
footer(s, num() + 1)

# ════════════════ 20. خارطة الطريق ════════════════
s = content_slide("خارطة الطريق — ما بعد المناقشة")
card(s, Inches(6.9), Inches(1.8), Inches(5.7), Inches(4.4),
     "قصير المدى (جاهز للتنفيذ)",
     "خط CI/CD على GitHub Actions: تحليل + اختبارات لكل تعديل\n\n"
     "مراقبة أعطال إنتاجية عبر Sentry بسياق المستخدم\n\n"
     "زر الطوارئ ومشاركة الرحلة مع جهة اتصال (SAF-01/02)\n\n"
     "قياسات أداء موثقة: إقلاع ≤ 3 ثوانٍ و 60 إطاراً/ثانية",
     accent=ORANGE, bsize=13)
card(s, Inches(0.7), Inches(1.8), Inches(5.7), Inches(4.4),
     "متوسط المدى (رؤية النمو)",
     "تتبع مباشر للرحلة على الخريطة أثناء تنفيذها\n\n"
     "ربط بوابة دفع إلكترونية عند توفرها محلياً\n\n"
     "نشر iOS على App Store + دعم إمكانية الوصول\n\n"
     "رحلات دورية (يومية للموظفين والطلاب) وتسعير ديناميكي",
     accent=CYAN, bsize=13)
footer(s, num() + 1)

# ════════════════ 21. الخاتمة ════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, DARK_NAVY)
add_rect(s, 0, 0, SW, Pt(6), ORANGE)
add_rect(s, 0, SH - Pt(6), SW, Pt(6), ORANGE)
logo_badge(s, SW / 2 - Inches(0.95), Inches(1.0), Inches(1.9))
add_text(s, Inches(1.5), Inches(3.1), Inches(10.33), Inches(0.9),
         [("شكراً لحسن استماعكم", 40, WHITE, True)], align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.15), Inches(10.33), Inches(0.6),
         [("عطريقك — لأن الطريق أوفر حين نتشاركه", 18, ORANGE, False)],
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(5.2), Inches(10.33), Inches(0.6),
         [("أسئلتكم موضع ترحيب", 16, CYAN, False)], align=PP_ALIGN.CENTER)

OUT = r"d:\my_projects\flutter\alatarekak\docs\Atariqak_Final_Presentation.pptx"
prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
